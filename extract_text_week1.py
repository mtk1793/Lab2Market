import os
from pypdf import PdfReader
from pptx import Presentation

week1_path = r"c:/Users/mtk17/OneDrive - Dalhousie University/Google Drive/PhD/Lab2Market/Presentations/Week 1"

def extract_from_pdf(path):
    try:
        reader = PdfReader(path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_from_pptx(path):
    try:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {e}"

print("Reading files from:", week1_path)
for filename in os.listdir(week1_path):
    file_path = os.path.join(week1_path, filename)
    
    # Skip the .md file
    if filename.endswith('.md') or filename.endswith('.resolved'):
        continue
        
    print(f"\n--- Processing {filename} ---")
    
    content = ""
    if filename.lower().endswith('.pdf'):
        content = extract_from_pdf(file_path)
    elif filename.lower().endswith('.pptx'):
        content = extract_from_pptx(file_path)
    else:
        print(f"Skipping unsupported file type: {filename}")
        continue

    with open("week1_content.txt", "a", encoding="utf-8") as f:
        f.write(f"\n\n--- FILE: {filename} ---\n\n")
        f.write(content)
    print(f"Extracted content from {filename}")

print("\nExtraction complete! Content saved to week1_content.txt")
