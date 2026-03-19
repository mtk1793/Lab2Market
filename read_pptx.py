from pptx import Presentation
from pptx.util import Inches, Pt
from lxml import etree
import os

pptx_path = r'c:\Users\mtk17\OneDrive\Desktop\L2M\Lab2Market\temp_presentation.pptx'

prs = Presentation(pptx_path)
print(f"Total Slides: {len(prs.slides)}\n")

# Focus on slides 5-9
for i in range(4, 9):  # 0-indexed, so slides 5-9
    slide = prs.slides[i]
    print(f"\n{'='*80}")
    print(f"SLIDE {i+1}")
    print('='*80)
    # Get all text from all shapes including in groups
    def get_text_from_shape(shape, indent=0):
        prefix = "  " * indent
        if hasattr(shape, 'text') and shape.text.strip():
            print(f"{prefix}[{shape.name}]: {shape.text.strip()!r}")
        if shape.shape_type == 6:  # Group shape
            for s in shape.shapes:
                get_text_from_shape(s, indent+1)
        if hasattr(shape, 'has_table') and shape.has_table:
            table = shape.table
            print(f"{prefix}[TABLE {shape.name}]:")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(f"{prefix}  {cells}")
    
    for shape in slide.shapes:
        get_text_from_shape(shape)
    
    # Also get the raw XML text content
    all_text = []
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for t in slide._element.iter(f'{{{ns}}}t'):
        if t.text and t.text.strip():
            all_text.append(t.text.strip())
    if all_text:
        print(f"\n  Raw text elements: {all_text}")

